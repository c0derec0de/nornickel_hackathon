import os
from ..config import settings

SUPPORTED = {".pdf", ".txt", ".md", ".docx", ".docm", ".pptx"}


def _read_pdf(path: str) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        return "\n".join((p.extract_text() or "") for p in reader.pages)
    except Exception:
        return ""


def _read_docx(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(c.text for c in row.cells))
        return "\n".join(parts)
    except Exception:
        return ""


def _read_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        return "\n".join(p for p in parts if p.strip())
    except Exception:
        return ""


def _read_file(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return _read_pdf(path)
    if ext in (".docx", ".docm"):
        return _read_docx(path)
    if ext == ".pptx":
        return _read_pptx(path)
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return ""


def chunk_text(text: str, size: int, overlap: int, max_chunks: int | None = None) -> list[str]:
    text = " ".join(text.split())
    if not text:
        return []
    chunks, start = [], 0
    while start < len(text):
        end = start + size
        chunks.append(text[start:end])
        start = end - overlap
        if max_chunks and len(chunks) >= max_chunks:
            break
    return chunks


def iter_files(corpus_dir: str, name_filter: list[str] | None = None):
    """Рекурсивно обходит корпус, возвращает пути поддерживаемых файлов."""
    for root, _dirs, files in os.walk(corpus_dir):
        for fname in sorted(files):
            if fname.startswith("."):
                continue
            if os.path.splitext(fname)[1].lower() not in SUPPORTED:
                continue
            if name_filter and not any(s.lower() in fname.lower() for s in name_filter):
                continue
            yield os.path.join(root, fname)


def load_corpus(corpus_dir: str | None = None, max_files: int | None = None,
                name_filter: list[str] | None = None, max_chunks_per_doc: int | None = None):
    """Возвращает список чанков: {doc, chunk_id, text}. Рекурсивно, мультиформатно.

    max_files          — ограничить число документов (контроль стоимости/времени)
    name_filter        — брать только файлы, чьё имя содержит одну из подстрок
    max_chunks_per_doc — не резать документ на более чем N чанков
    """
    corpus_dir = corpus_dir or settings.corpus_dir
    items = []
    if not os.path.isdir(corpus_dir):
        return items

    count = 0
    for path in iter_files(corpus_dir, name_filter=name_filter):
        if max_files and count >= max_files:
            break
        text = _read_file(path)
        chunks = chunk_text(text, settings.chunk_size, settings.chunk_overlap,
                            max_chunks=max_chunks_per_doc)
        if not chunks:
            continue
        doc = os.path.relpath(path, corpus_dir)
        for i, ch in enumerate(chunks):
            items.append({"doc": doc, "chunk_id": f"{doc}::{i}", "text": ch})
        count += 1
    return items
